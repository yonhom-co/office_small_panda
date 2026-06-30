"""gen_ppt 工具 —— 基于 python-pptx 生成 PPT 交付物。

模型在分析完成后调用，传标题、幻灯片列表（每页含标题/要点/可选图表路径），
生成 .pptx 文件。复刻「端到端产出 PPT」。
"""
from __future__ import annotations

import os
import uuid
from pathlib import Path

from .base import Tool, default_registry
from ..dataset import WORKDIR


def _gen_ppt(params: dict, shared: dict) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    title = params.get("title", "数据分析汇报")
    slides = params.get("slides", [])
    charts = shared.get("charts", [])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True

    chart_idx = 0
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 标题
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = s.get("heading", "")
        p.font.size = Pt(28)
        p.font.bold = True
        # 要点
        bullets = s.get("bullets", [])
        if bullets:
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(7.5), Inches(5.5))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = f"• {b}"
                para.font.size = Pt(18)
        # 图表（每页一张，从 shared charts 顺序取）
        if s.get("with_chart") and chart_idx < len(charts):
            try:
                slide.shapes.add_picture(charts[chart_idx], Inches(8.2), Inches(1.6),
                                         width=Inches(4.8))
                chart_idx += 1
            except Exception:
                pass

    out_dir = WORKDIR / "reports"
    out_dir.mkdir(parents=True, exist_ok=True)
    name = f"report_{uuid.uuid4().hex[:8]}.pptx"
    out_path = out_dir / name
    prs.save(str(out_path))
    shared["ppt_path"] = str(out_path)
    return (f"PPT 已生成：{out_path}\n（共 {len(slides) + 1} 页，"
            f"含 {chart_idx} 张图表）")


gen_ppt_tool = Tool(
    name="gen_ppt",
    description=(
        "生成 PPT 汇报文件。分析完成后调用。"
        "传 title（标题）、slides（幻灯片列表，每项含 heading 标题、bullets 要点列表、"
        "可选 with_chart 是否嵌入图表）。已生成的图表按顺序嵌入。"
    ),
    input_schema={
        "type": "object",
        "properties": {
            "title": {"type": "string"},
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {"type": "string"},
                        "bullets": {"type": "array", "items": {"type": "string"}},
                        "with_chart": {"type": "boolean"},
                    },
                    "required": ["heading"],
                },
            },
        },
        "required": ["title"],
    },
    executor=_gen_ppt,
)

default_registry.register(gen_ppt_tool)
